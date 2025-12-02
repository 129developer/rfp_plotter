"""
PowerPoint generation utilities for creating RFP proposal presentations.
Converts JSON proposal data into professional PowerPoint presentations.
"""

import os
import tempfile
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
from datetime import datetime
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

from ..models.rfp_models import RFPProposal, Phase, CostItem, Milestone, UserStory

logger = logging.getLogger(__name__)


class PPTGenerator:
    """Utility class for generating PowerPoint presentations from RFP proposals"""
    
    def __init__(self):
        """Initialize the PPT generator"""
        self.presentation = None
        self.slide_layouts = None
        
        # Define color scheme
        self.colors = {
            'primary': RGBColor(0x1f, 0x4e, 0x79),      # Dark blue
            'secondary': RGBColor(0x2e, 0x86, 0xab),    # Light blue
            'accent': RGBColor(0xf2, 0x8e, 0x2b),       # Orange
            'text': RGBColor(0x2c, 0x2c, 0x2c),         # Dark gray
            'light_gray': RGBColor(0xf5, 0xf5, 0xf5),   # Light gray
        }
    
    def generate_presentation(self, proposal: RFPProposal, output_path: Optional[str] = None) -> str:
        """
        Generate a complete PowerPoint presentation from an RFP proposal.
        
        Args:
            proposal: RFP proposal data
            output_path: Optional output file path
            
        Returns:
            Path to the generated PowerPoint file
        """
        try:
            logger.info("Starting PowerPoint generation...")
            
            # Create new presentation
            self.presentation = Presentation()
            self.slide_layouts = self.presentation.slide_layouts
            
            # Generate slides
            self._create_cover_slide(proposal.cover)
            self._create_meta_slide(proposal.meta)
            self._create_background_objectives_slide(proposal.background_and_objectives)
            
            # Create phase slides
            for phase in proposal.phases:
                self._create_phase_slide(phase)
                self._create_services_slide(phase)
                self._create_assumptions_slide(phase)
                self._create_dependencies_slide(phase)
            
            # Create architecture and deployment slides
            self._create_architecture_slide(proposal.solution_architecture)
            self._create_deployment_slide(proposal.deployment_view)
            
            # Create plan and commercials slides
            self._create_plan_slide(proposal.plan)
            self._create_commercials_slide(proposal.commercials)
            
            # Create appendix slides
            if proposal.appendix_stories.user_stories:
                self._create_user_stories_slide(proposal.appendix_stories.user_stories)
            
            if proposal.components:
                self._create_components_slide(proposal.components)
            
            # Save presentation
            if not output_path:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                client_name = proposal.cover.client_name.replace(' ', '_') if proposal.cover.client_name else "Client"
                output_path = f"{client_name}_Proposal_{timestamp}.pptx"
            
            self.presentation.save(output_path)
            logger.info(f"PowerPoint presentation saved to: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint presentation: {e}")
            raise
    
    def _create_cover_slide(self, cover_info):
        """Create the cover slide"""
        slide_layout = self.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = cover_info.project_title or "RFP Proposal"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle_text = f"Proposal for {cover_info.client_name}\n"
        subtitle_text += f"Version {cover_info.version} | {cover_info.date}\n"
        subtitle_text += f"Prepared by {cover_info.vendor_name}"
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['secondary']
    
    def _create_meta_slide(self, meta_info):
        """Create the metadata slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Document Information"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add metadata
        p = text_frame.paragraphs[0]
        p.text = f"Document Version: {meta_info.doc_version}"
        p.font.size = Pt(18)
        
        p = text_frame.add_paragraph()
        p.text = f"Prepared by: {meta_info.prepared_by}"
        p.font.size = Pt(18)
        
        p = text_frame.add_paragraph()
        p.text = f"Contact: {meta_info.contact_email}"
        p.font.size = Pt(18)
        
        p = text_frame.add_paragraph()
        p.text = ""
        
        p = text_frame.add_paragraph()
        p.text = meta_info.confidentiality_note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = self.colors['text']
    
    def _create_background_objectives_slide(self, background_objectives):
        """Create the background and objectives slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Background & Objectives"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add context
        if background_objectives.context:
            p = text_frame.paragraphs[0]
            p.text = "Context"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            p = text_frame.add_paragraph()
            p.text = background_objectives.context
            p.font.size = Pt(16)
            p.level = 1
        
        # Add objectives
        if background_objectives.key_objectives:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Key Objectives"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for objective in background_objectives.key_objectives:
                p = text_frame.add_paragraph()
                p.text = objective
                p.font.size = Pt(16)
                p.level = 1
        
        # Add current state
        if background_objectives.current_state:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Current State"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            p = text_frame.add_paragraph()
            p.text = background_objectives.current_state
            p.font.size = Pt(16)
            p.level = 1
    
    def _create_phase_slide(self, phase: Phase):
        """Create a slide for a project phase"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = f"{phase.title} - Scope"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add scope summary
        if phase.scope_summary:
            p = text_frame.paragraphs[0]
            p.text = phase.scope_summary
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['text']
        
        # Add deliverables
        if phase.deliverables:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Deliverables"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for deliverable in phase.deliverables:
                p = text_frame.add_paragraph()
                p.text = deliverable
                p.font.size = Pt(16)
                p.level = 1
        
        # Add acceptance criteria
        if phase.acceptance_criteria:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Acceptance Criteria"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for criteria in phase.acceptance_criteria:
                p = text_frame.add_paragraph()
                p.text = criteria
                p.font.size = Pt(16)
                p.level = 1
    
    def _create_services_slide(self, phase: Phase):
        """Create a services slide for a phase"""
        if not phase.services.service_list:
            return
        
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = f"{phase.title} - Services"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = "Services Included"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        for service in phase.services.service_list:
            p = text_frame.add_paragraph()
            p.text = service
            p.font.size = Pt(18)
            p.level = 1
            
            # Add description if available
            if service in phase.services.service_descriptions:
                p = text_frame.add_paragraph()
                p.text = phase.services.service_descriptions[service]
                p.font.size = Pt(14)
                p.level = 2
                p.font.color.rgb = self.colors['text']
    
    def _create_assumptions_slide(self, phase: Phase):
        """Create an assumptions slide for a phase"""
        if not phase.assumptions:
            return
        
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = f"{phase.title} - Assumptions"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = "Key Assumptions"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        for assumption in phase.assumptions:
            p = text_frame.add_paragraph()
            p.text = assumption
            p.font.size = Pt(16)
            p.level = 1
    
    def _create_dependencies_slide(self, phase: Phase):
        """Create a dependencies slide for a phase"""
        if not phase.dependencies:
            return
        
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = f"{phase.title} - Dependencies"
        
        # Create table for dependencies
        rows = len(phase.dependencies) + 1  # +1 for header
        cols = 3
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.5 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set header
        table.cell(0, 0).text = "Dependency"
        table.cell(0, 1).text = "Owner"
        table.cell(0, 2).text = "Lead Time"
        
        # Style header
        for col in range(cols):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Add dependency data
        for i, dep in enumerate(phase.dependencies, 1):
            table.cell(i, 0).text = dep.dependency
            table.cell(i, 1).text = dep.owner
            table.cell(i, 2).text = dep.lead_time
    
    def _create_architecture_slide(self, architecture):
        """Create the solution architecture slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Solution Architecture"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add architecture summary
        if architecture.architecture_summary:
            p = text_frame.paragraphs[0]
            p.text = architecture.architecture_summary
            p.font.size = Pt(16)
        
        # Add technology choices
        if architecture.key_technology_choices:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Key Technology Choices"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for choice in architecture.key_technology_choices:
                p = text_frame.add_paragraph()
                p.text = choice
                p.font.size = Pt(16)
                p.level = 1
        
        # Add diagram placeholder
        if architecture.diagram_spec.nodes:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Architecture Components"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for node in architecture.diagram_spec.nodes:
                p = text_frame.add_paragraph()
                p.text = node
                p.font.size = Pt(14)
                p.level = 1
    
    def _create_deployment_slide(self, deployment):
        """Create the deployment view slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Deployment Architecture"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add environments
        if deployment.environments:
            p = text_frame.paragraphs[0]
            p.text = "Deployment Environments"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for env in deployment.environments:
                p = text_frame.add_paragraph()
                p.text = env
                p.font.size = Pt(16)
                p.level = 1
        
        # Add networking notes
        if deployment.networking_notes:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Networking Considerations"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            p = text_frame.add_paragraph()
            p.text = deployment.networking_notes
            p.font.size = Pt(16)
            p.level = 1
    
    def _create_plan_slide(self, plan):
        """Create the project plan slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Project Plan"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add methodology
        p = text_frame.paragraphs[0]
        p.text = f"Methodology: {plan.methodology}"
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"Sprint Length: {plan.sprint_length_days} days"
        p.font.size = Pt(16)
        
        # Add milestones
        if plan.milestones:
            p = text_frame.add_paragraph()
            p.text = ""
            
            p = text_frame.add_paragraph()
            p.text = "Key Milestones"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            for milestone in plan.milestones:
                p = text_frame.add_paragraph()
                p.text = f"{milestone.name} - {milestone.date}"
                p.font.size = Pt(16)
                p.level = 1
                
                if milestone.description:
                    p = text_frame.add_paragraph()
                    p.text = milestone.description
                    p.font.size = Pt(14)
                    p.level = 2
                    p.font.color.rgb = self.colors['text']
    
    def _create_commercials_slide(self, commercials):
        """Create the commercials slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Commercial Proposal"
        
        # Create cost table if available
        if commercials.cost_table:
            rows = len(commercials.cost_table) + 2  # +1 for header, +1 for total
            cols = 3
            
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.5 * rows)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set header
            table.cell(0, 0).text = "Item"
            table.cell(0, 1).text = "Description"
            table.cell(0, 2).text = "Cost"
            
            # Style header
            for col in range(cols):
                cell = table.cell(0, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.colors['primary']
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True
            
            # Add cost data
            total_cost = 0
            for i, cost_item in enumerate(commercials.cost_table, 1):
                table.cell(i, 0).text = cost_item.item
                table.cell(i, 1).text = cost_item.description
                table.cell(i, 2).text = f"${cost_item.cost:,.2f}"
                total_cost += cost_item.cost
            
            # Add total row
            total_row = len(commercials.cost_table) + 1
            table.cell(total_row, 0).text = "TOTAL"
            table.cell(total_row, 1).text = ""
            table.cell(total_row, 2).text = f"${total_cost:,.2f}"
            
            # Style total row
            for col in range(cols):
                cell = table.cell(total_row, col)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.colors['light_gray']
        
        # Add payment terms
        if commercials.payment_terms:
            # Add text box for payment terms
            left = Inches(1)
            top = Inches(6)
            width = Inches(8)
            height = Inches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            
            p = text_frame.paragraphs[0]
            p.text = f"Payment Terms: {commercials.payment_terms}"
            p.font.size = Pt(16)
            p.font.bold = True
    
    def _create_user_stories_slide(self, user_stories: List[UserStory]):
        """Create user stories slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "User Stories"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, story in enumerate(user_stories[:5], 1):  # Limit to 5 stories
            p = text_frame.paragraphs[0] if i == 1 else text_frame.add_paragraph()
            p.text = f"Story {i}: As a {story.role}, I want {story.goal} so that {story.benefit}"
            p.font.size = Pt(14)
            p.level = 0
            
            if story.acceptance_criteria:
                for criteria in story.acceptance_criteria[:2]:  # Limit to 2 criteria per story
                    p = text_frame.add_paragraph()
                    p.text = f"• {criteria}"
                    p.font.size = Pt(12)
                    p.level = 1
                    p.font.color.rgb = self.colors['text']
            
            if i < len(user_stories[:5]):  # Add spacing between stories
                text_frame.add_paragraph()
    
    def _create_components_slide(self, components):
        """Create technology components slide"""
        slide_layout = self.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Technology Components"
        
        # Add content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, component in enumerate(components[:6], 1):  # Limit to 6 components
            p = text_frame.paragraphs[0] if i == 1 else text_frame.add_paragraph()
            p.text = component.name
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['primary']
            
            p = text_frame.add_paragraph()
            p.text = component.description
            p.font.size = Pt(14)
            p.level = 1
            
            if component.rationale:
                p = text_frame.add_paragraph()
                p.text = f"Rationale: {component.rationale}"
                p.font.size = Pt(12)
                p.level = 1
                p.font.color.rgb = self.colors['text']
            
            if i < len(components[:6]):  # Add spacing
                text_frame.add_paragraph()


def generate_ppt_from_proposal(proposal: RFPProposal, output_path: Optional[str] = None) -> str:
    """
    Convenience function to generate PowerPoint from RFP proposal.
    
    Args:
        proposal: RFP proposal data
        output_path: Optional output file path
        
    Returns:
        Path to generated PowerPoint file
    """
    generator = PPTGenerator()
    return generator.generate_presentation(proposal, output_path)